# output_generator.py
import os
import uuid
import logging
from datetime import datetime
from typing import List, Dict, Optional

import pandas as pd
from flask import current_app, url_for

logger = logging.getLogger(__name__)


class OutputGenerator:
    """多格式输出生成器：将结果数据集导出为 HTML/Excel/Word/PPT 文件"""

    def __init__(self, output_dir: Optional[str] = None):
        self.output_dir = output_dir or current_app.config['OUTPUT_FOLDER']
        os.makedirs(self.output_dir, exist_ok=True)

    def generate(self,
                 data: List[Dict],
                 format: str = 'excel',
                 context: str = '',
                 chart_files: Optional[List[str]] = None) -> str:
        """
        生成导出文件
        :param data: 数据集，列表套字典
        :param format: 输出格式 (excel / html / word / ppt)
        :param context: 业务上下文（用于标题或说明）
        :param chart_files: 图表文件路径列表（可选，主要用于 HTML 嵌入）
        :return: 生成的文件名（不含路径前缀）
        """
        if not data:
            raise ValueError("数据集为空，无法生成文件")

        format = format.lower()
        if format not in ('excel', 'html', 'word', 'ppt'):
            raise ValueError(f"不支持的输出格式: {format}")

        # 生成唯一文件名
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        unique_id = str(uuid.uuid4())[:8]
        file_name = f"analysis_{timestamp}_{unique_id}"
        file_path = os.path.join(self.output_dir, file_name)

        try:
            if format == 'excel':
                file_path = self._to_excel(data, file_path, context)
            elif format == 'html':
                file_path = self._to_html(data, file_path, context, chart_files)
            elif format == 'word':
                file_path = self._to_word(data, file_path, context)
            elif format == 'ppt':
                file_path = self._to_ppt(data, file_path, context)

            # 记录日志
            logger.info(f"文件已生成: {file_path}")
            return os.path.basename(file_path)

        except Exception as e:
            logger.error(f"生成 {format} 文件失败: {str(e)}")
            raise RuntimeError(f"文件生成失败: {str(e)}")

    # ---------- Excel ----------
    def _to_excel(self, data: List[Dict], file_path: str, context: str) -> str:
        df = pd.DataFrame(data)
        # 限制数据行数，避免文件过大（可配置）
        max_rows = 10000
        if len(df) > max_rows:
            df = df.head(max_rows)
            logger.warning(f"数据集过大，已截取前 {max_rows} 行")
        full_path = file_path + '.xlsx'
        with pd.ExcelWriter(full_path, engine='openpyxl') as writer:
            df.to_excel(writer, index=False, sheet_name='分析结果')
            # 添加说明工作表
            info_df = pd.DataFrame({'生成时间': [datetime.now().strftime('%Y-%m-%d %H:%M:%S')],
                                    '上下文': [context]})
            info_df.to_excel(writer, index=False, sheet_name='说明')
        return full_path

    # ---------- HTML ----------
    def _to_html(self, data: List[Dict], file_path: str, context: str,
                 chart_files: Optional[List[str]] = None) -> str:
        df = pd.DataFrame(data)
        # 使用 Jinja2 模板渲染
        template_str = """<!DOCTYPE html>
<html lang="zh">
<head>
    <meta charset="UTF-8">
    <title>数据分析结果</title>
    <style>
        body { font-family: 'Microsoft YaHei', sans-serif; margin: 20px; }
        h1 { color: #333; }
        table { border-collapse: collapse; width: 100%; }
        th, td { border: 1px solid #ddd; padding: 8px; text-align: left; }
        th { background-color: #4CAF50; color: white; }
        tr:nth-child(even) { background-color: #f2f2f2; }
        .chart { margin-top: 20px; }
    </style>
</head>
<body>
    <h1>数据分析结果</h1>
    {% if context %}<p><strong>分析背景：</strong>{{ context }}</p>{% endif %}
    {{ table_html | safe }}
    {% if charts %}
    <div class="chart">
        <h2>图表</h2>
        {% for chart in charts %}
        <img src="{{ chart }}" alt="图表" style="max-width:100%; margin-bottom:10px;">
        {% endfor %}
    </div>
    {% endif %}
</body>
</html>"""

        from jinja2 import Template
        template = Template(template_str)
        table_html = df.to_html(classes='data', index=False, border=0, justify='left')
        # 处理图表文件：将其转换为相对于 static 的 URL
        chart_urls = []
        if chart_files:
            for f in chart_files:
                # f 是绝对路径，需转换为 /static/outputs/... 的相对链接
                fname = os.path.basename(f)
                # 将图表文件拷贝到输出目录（如果不在同一目录）
                if os.path.dirname(f) != self.output_dir:
                    import shutil
                    dest = os.path.join(self.output_dir, fname)
                    shutil.copy2(f, dest)
                    logger.debug(f"图表已复制到 {dest}")
                chart_urls.append(url_for('static', filename=f'outputs/{fname}'))

        html_content = template.render(context=context, table_html=table_html, charts=chart_urls)

        full_path = file_path + '.html'
        with open(full_path, 'w', encoding='utf-8') as f:
            f.write(html_content)
        return full_path

    # ---------- Word ----------
    def _to_word(self, data: List[Dict], file_path: str, context: str) -> str:
        try:
            from docx import Document
        except ImportError:
            raise ImportError("python-docx 未安装，无法生成 Word 文档")

        doc = Document()
        doc.add_heading('数据分析报告', 0)
        if context:
            doc.add_paragraph(f'分析背景：{context}')

        # 限制行数，避免文档过大
        max_rows = 500
        if len(data) > max_rows:
            data = data[:max_rows]
            logger.warning(f"数据集过大，仅导出前 {max_rows} 行到 Word")

        if data:
            table = doc.add_table(rows=1, cols=len(data[0]))
            table.style = 'Light Shading Accent 1'
            # 表头
            hdr_cells = table.rows[0].cells
            for i, key in enumerate(data[0].keys()):
                hdr_cells[i].text = str(key)
            # 数据行
            for row_dict in data:
                row_cells = table.add_row().cells
                for i, val in enumerate(row_dict.values()):
                    row_cells[i].text = str(val) if val is not None else ''

        doc.add_paragraph(f'生成时间：{datetime.now().strftime("%Y-%m-%d %H:%M:%S")}')
        full_path = file_path + '.docx'
        doc.save(full_path)
        return full_path

    # ---------- PPT ----------
    def _to_ppt(self, data: List[Dict], file_path: str, context: str) -> str:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            raise ImportError("python-pptx 未安装，无法生成 PPT 文件")

        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # 标题和内容布局
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "数据分析结果"
        if context:
            subtitle = slide.placeholders[1]
            subtitle.text = context

        # 添加表格到新幻灯片
        max_rows = 30
        if len(data) > max_rows:
            data = data[:max_rows]
            logger.warning(f"数据集过大，仅导出前 {max_rows} 行到 PPT")

        if data:
            rows = len(data) + 1
            cols = len(data[0])
            table_slide = prs.slides.add_slide(prs.slide_layouts[5])  # 空白
            table_shape = table_slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(9), Inches(5))
            table = table_shape.table

            # 表头
            for i, key in enumerate(data[0].keys()):
                cell = table.cell(0, i)
                cell.text = str(key)
                # 设置表头字体
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(10)
                    paragraph.font.bold = True

            # 数据
            for r, row_dict in enumerate(data):
                for c, val in enumerate(row_dict.values()):
                    cell = table.cell(r+1, c)
                    cell.text = str(val) if val is not None else ''
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(9)

        full_path = file_path + '.pptx'
        prs.save(full_path)
        return full_path

    # ---------- 清理过期文件 ----------
    @staticmethod
    def cleanup_old_files(hours: int = 1):
        """删除超过指定小时数的输出文件，可在定时任务中调用"""
        import time
        output_dir = current_app.config['OUTPUT_FOLDER']
        now = time.time()
        for fname in os.listdir(output_dir):
            fpath = os.path.join(output_dir, fname)
            if os.path.isfile(fpath) and (now - os.path.getmtime(fpath)) > hours * 3600:
                try:
                    os.remove(fpath)
                    logger.info(f"清理过期文件: {fname}")
                except Exception as e:
                    logger.warning(f"清理文件失败 {fname}: {e}")